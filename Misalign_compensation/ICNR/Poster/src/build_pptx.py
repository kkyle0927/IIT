#!/usr/bin/env python3
"""Assemble the ICNR 2026 poster (A0 portrait) as an editable .pptx.

Five macro sections — Background / Robot System / Methodology / Experiment /
Results. Section titles are plain text over a rule, not filled bars, and carry
no numbering. Every synthesis sentence lives inside that section's own bullet
list — there is no separate "takeaway" callout box anywhere on the poster
(see the design-rule note below), and the footer is a single disclaimer line.

Layout is modular: PAGE_W / PAGE_H / MARGIN / GUTTER set the frame, and each
section owns one height constant (SEC_*). A different poster size needs only
those constants changed. All text, boxes and arrows are native PowerPoint
objects; only the plots, the architecture diagram and the two logos are
images.

Design rule this file follows throughout (generalizable to any AI-authored
slide/poster): no grey text — de-emphasise with size/weight, never with colour,
because low-contrast grey is a routine AI-generated-slide tell and it is also
worse for a poster read from 1-2 m away; and no colour-block "label" accents
(a tinted rectangle or a coloured bar glued to the edge of a box to mark it as
"the important one") — a plain full-width structural rule (e.g. under a
section title) is fine, but a rectangle whose only job is to flag one specific
piece of content is not used here.

Content comes from IJCAS/revision_submit/manuscript.tex and the locked evidence
under Experiment/paper_evidence/. No value here is estimated.
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

ASSETS = Path(__file__).resolve().parents[1] / "assets"
OUT = Path(__file__).resolve().parents[1] / "ICNR_Poster.pptx"

# ------------------------------------------------------------------- LAYOUT --
PAGE_W, PAGE_H = 84.1, 118.9          # A0 portrait, cm
MARGIN, GUTTER, GAP = 3.5, 2.6, 1.1
CONTENT_W = PAGE_W - 2 * MARGIN
COL_W = (CONTENT_W - GUTTER) / 2
COL_L, COL_R = MARGIN, MARGIN + COL_W + GUTTER
PAD = 0.5
INNER_W = COL_W - 2 * PAD

HEADER_H, FOOT_H = 15.0, 2.2
BODY_Y = MARGIN + HEADER_H + GAP
FOOT_Y = PAGE_H - MARGIN - FOOT_H

SEC_HEAD_H = 3.6                      # section title + rule
RES_GAP_TITLE_IMG, RES_GAP_IMG_TXT, RES_GAP_BLOCK = 0.05, 0.08, 0.1

# --------------------------------------------------------------------- STYLE --
BLUE = RGBColor(0x1F, 0x48, 0x99)      # KAIST blue
GREEN = RGBColor(0x9F, 0xBE, 0x19)     # EXO LAB green — used only for the lab wordmark
DARK = RGBColor(0x3D, 0x3F, 0x36)      # EXO LAB dark
INK = RGBColor(0x1A, 0x1A, 0x1A)       # all body/caption text — no grey anywhere
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_EDGE = RGBColor(0xD5, 0xDB, 0xE3)
SOFT = RGBColor(0xF3, 0xF6, 0xFA)
FONT = "Times New Roman"

T_SECTION, T_SUB, T_BODY, T_CHIP = 52, 32, 26, 26


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.5):
    s = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.text_frame.text = ""
    return s


def text(slide, x, y, w, h, runs, size=T_BODY, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6, line=0.95, italic=False):
    """runs: a string; a (string, overrides) paragraph; or a list of those.

    A paragraph item may itself be a *list* of (substring, overrides) run
    tuples, for inline mixed formatting within one line/paragraph — used for
    proper subscripts (pass overrides={"sub": True} on the subscripted run).
    """
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate([runs] if isinstance(runs, str) else runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, str):
            segs, para_over = [(item, {})], {}
        elif isinstance(item, tuple):
            segs, para_over = [item], item[1]
        else:
            segs, para_over = item, {}
        p.alignment = para_over.get("align", align)
        p.space_after = Pt(para_over.get("space", space))
        p.line_spacing = para_over.get("line", line)
        for seg_text, seg_over in segs:
            r = p.add_run()
            r.text = seg_text
            f = r.font
            f.name = FONT
            base_size = seg_over.get("size", para_over.get("size", size))
            if seg_over.get("sub"):
                f.size = Pt(round(base_size * 0.66))
                r._r.get_or_add_rPr().set("baseline", "-25000")
            else:
                f.size = Pt(base_size)
            f.bold = seg_over.get("bold", para_over.get("bold", bold))
            f.italic = seg_over.get("italic", para_over.get("italic", italic))
            f.color.rgb = seg_over.get("color", para_over.get("color", color))
    return box


def image(slide, path, x, y, w, h):
    """Place an image scaled to fit inside (w, h) and centred."""
    iw, ih = Image.open(ASSETS / path).size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    return slide.shapes.add_picture(str(ASSETS / path), Cm(x + (w - dw) / 2),
                                    Cm(y + (h - dh) / 2), Cm(dw), Cm(dh))


def chart_h(fname, width=INNER_W):
    """Placed height of a poster-column-width chart, from its saved aspect ratio."""
    iw, ih = Image.open(ASSETS / fname).size
    return width * ih / iw


def section(slide, x, y, w, title):
    """Section title over a rule — open layout, no card border or fill.
    Sections are separated by whitespace (GAP) and this rule, not by a box."""
    text(slide, x + PAD, y + 0.3, w - 2 * PAD, 2.4, title, size=T_SECTION,
         color=BLUE, bold=True)
    rect(slide, x + PAD, y + SEC_HEAD_H - 0.55, w - 2 * PAD, 0.16, fill=BLUE)
    return y + SEC_HEAD_H


def sub(slide, x, y, w, title):
    text(slide, x, y, w, 1.7, title, size=T_SUB, color=DARK, bold=True)
    return y + 1.7


def bullets(slide, x, y, w, items, size=T_BODY, space=9, h=99, anchor=MSO_ANCHOR.TOP):
    """A plain bullet list — every bullet is "• " + a string or run-list.
    Pass h (the row's allocated height) with anchor=MIDDLE to vertically
    centre the whole list within that height, matching an adjacent image."""
    paras = []
    for it in items:
        if isinstance(it, str):
            paras.append("•  " + it)
        else:  # list of (text, overrides) runs: prefix the bullet on the first
            paras.append([("•  " + it[0][0], it[0][1])] + list(it[1:]))
    text(slide, x, y, w, h, paras, size=size, color=INK, space=space, line=0.95,
         anchor=anchor)


PT_CM = 0.03528
CHAR_W_EM = 0.44    # Times New Roman average glyph advance, in em
LEADING = 1.22       # line height as a multiple of font size (calibrated against render)


def _item_len(it):
    return len(it) if isinstance(it, str) else sum(len(seg[0]) for seg in it)


def text_h(items, width, size=T_BODY, space=9, bulleted=True):
    """Predicted height (cm) of `items` laid out by text()/bullets() at `width`
    and `size` — used to size a section from its actual content instead of a
    hand-guessed constant. Calibrated empirically against the LibreOffice
    render; deliberately generous (LEADING, CHAR_W_EM) so a mis-estimate errs
    toward extra whitespace, never overlap."""
    if isinstance(items, str):
        items = [items]
    prefix = 3 if bulleted else 0  # "•  "
    cpl = max(10, width / (size * PT_CM * CHAR_W_EM))
    lines = sum(max(1, -(-(prefix + _item_len(it)) // int(cpl))) for it in items)
    return lines * size * PT_CM * LEADING + max(0, len(items) - 1) * space * PT_CM


def chip(slide, x, y, w, h, body, fill=SOFT, color=BLUE, size=T_CHIP, bold=True):
    rect(slide, x, y, w, h, fill=fill, line=CARD_EDGE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x + 0.3, y, w - 0.6, h, body, size=size, bold=bold, color=color,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=0.95)


# ===========================================================================  =
prs = Presentation()
prs.slide_width, prs.slide_height = Cm(PAGE_W), Cm(PAGE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
rect(slide, 0, 0, PAGE_W, PAGE_H, fill=WHITE)

# ------------------------------------------------------------------- header --
rect(slide, 0, 0, PAGE_W, MARGIN + HEADER_H, fill=RGBColor(0xE9, 0xED, 0xF5))
rect(slide, 0, MARGIN + HEADER_H - 0.30, PAGE_W, 0.30, fill=BLUE)

image(slide, "logo_kaist.png", MARGIN, MARGIN + 0.7, 11.5, 3.4)
ex = image(slide, "logo_exolab.png", PAGE_W - MARGIN - 9.0, MARGIN + 0.1, 9.0, 4.0)
exl, exw = ex.left / 360000, ex.width / 360000
text(slide, exl - 2.0, MARGIN + 4.2, exw + 4.0, 1.0, "EXOSKELETON LABORATORY",
     size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
text(slide, MARGIN, MARGIN + 4.4, 20.0, 1.2, "ICNR 2026", size=26, bold=True, color=INK)

text(slide, MARGIN, MARGIN + 5.9, CONTENT_W, 4.4,
     "Data-Driven Compensation of Human–Robot Misalignment for "
     "Human Hip Kinematics Estimation in Wearable Robots",
     size=54, bold=True, color=BLUE, align=PP_ALIGN.CENTER, line=0.95)
text(slide, MARGIN, MARGIN + 9.9, CONTENT_W, 1.6,
     "Chanyoung Ko and Kyoungchul Kong*", size=34, bold=True, color=INK,
     align=PP_ALIGN.CENTER)
text(slide, MARGIN, MARGIN + 11.4, CONTENT_W, 1.4,
     "Exoskeleton Laboratory, Department of Mechanical Engineering, KAIST, Daejeon, "
     "Republic of Korea   ·   kkyle0927@kaist.ac.kr;  kckong@kaist.ac.kr  "
     "(*corresponding author)",
     size=23, color=INK, align=PP_ALIGN.CENTER)
text(slide, MARGIN, MARGIN + 12.9, CONTENT_W, 1.6,
     "The wearer's own thigh angle, read from the robot's own sensors. "
     "No body-worn sensors, no motion capture.",
     size=30, bold=True, color=BLUE, align=PP_ALIGN.CENTER, line=0.95)

def img_h(fname, width):
    iw, ih = Image.open(ASSETS / fname).size
    return width * ih / iw


BOTTOM_PAD = 0.3   # breathing room below the last element in every section


# A single shared image-slot width and text start column for every
# image-beside-bullets section (Background / Robot System / Experiment): each
# image keeps its own chosen display size but is centred within this common
# slot, so every section's body text starts at the same x — and, within each
# row, the image and the bullet list are both vertically centred on the same
# line so a shorter one never looks like it is "floating" against a taller one.
IMG_SLOT_W = 20.5
TEXT_X_OFFSET = IMG_SLOT_W + 2 * PAD
ROW_TEXT_W = COL_W - IMG_SLOT_W - 2 * PAD


# ---------------------------------------------------------- Background layout
BG_ITEMS = [
    "Robot-mounted sensors measure the robot, not the wearer's anatomical "
    "thigh angle.",
    "Trunk rotation skews the trunk-fixed robot frame off the "
    "anatomical hip axis.",
    "Body-worn IMUs would fix this, at a usability and form-factor cost.",
    "The raw robot angle is off by 6.72° on average, worse when "
    "assistance is on.",
]
BG_IMG_W = 19.5
BG_SIZE = T_BODY
BG_TEXT_X = COL_L + TEXT_X_OFFSET
BG_TEXT_W = ROW_TEXT_W
BG_ROW_H = max(img_h("fig_p1_concept.png", BG_IMG_W),
               text_h(BG_ITEMS, BG_TEXT_W, size=BG_SIZE, space=11))
SEC_BACKGROUND = SEC_HEAD_H + BG_ROW_H + BOTTOM_PAD

# --------------------------------------------------------- Robot System layout
RS_ITEMS = [
    "Bilateral hip exoskeleton (Angel KIT H10), 10 N·m per side",
    "Onboard sensing: hip encoders, trunk IMU, pelvis-side robot-frame IMU",
    [("θ", {}), ("robot", {"sub": True}), (" = θ", {}), ("encoder", {"sub": True}),
     (" − θ", {}), ("IMU", {"sub": True}), (" − θ", {}), ("calibration", {"sub": True}),
     (", from a T-pose standing calibration after donning", {})],
    [("τ", {}), ("peak", {"sub": True}),
     (" = 4 / 7 N·m (lv4 / lv7),  T", {}), ("p", {"sub": True}), (" = 0.1 s", {})],
    "Embedded module: NVIDIA Jetson Orin NX",
    "Everything the estimator needs is already mounted on the robot.",
]
RS_IMG_W = 18.2
RS_SIZE = T_BODY
RS_TEXT_X = COL_L + TEXT_X_OFFSET
RS_TEXT_W = ROW_TEXT_W
RS_ROW_H = max(img_h("fig_p3_sensors.png", RS_IMG_W),
               text_h(RS_ITEMS, RS_TEXT_W, size=RS_SIZE, space=9))
SEC_ROBOT = SEC_HEAD_H + RS_ROW_H + BOTTOM_PAD

# ----------------------------------------------------------- Methodology flow
BIOMECH_INTRO = ("Trunk and pelvis counter-rotate during gait, and the robot frame "
                 "rides on the trunk, so trunk yaw carries information about the "
                 "discrepancy that the robot's own IMU already measures.")
YAW_CAPTION = ("Relative yaw grows with walking speed and is negatively correlated "
               "with the discrepancy in all ten participants analysed.")
CALIB_NOTE = ("Calibration: one 1.0-s standing window after donning. No "
              "motion-capture label is used at any point during inference.")
PARAM_LINE = ("8,658 parameters  ·  fully causal  ·  no actuator current  ·  "
              "no body-worn sensor")
METHOD_CLOSE = ("A biomechanically motivated input, compact enough for the "
                "robot's own embedded computer.")
CHIP_H = 3.6


def layout_methodology(y0, draw):
    """Lay out the Methodology section top-to-bottom; returns the y just below
    the last element. Called once with draw=False to measure the required
    section height (so the section() box can be drawn at the right size),
    then again with draw=True at the final y to actually place everything —
    both passes share this one function so the two can never drift apart."""
    cur = y0
    if draw:
        cur = sub(slide, COL_L + PAD, cur, INNER_W, "Biomechanical cue")
    else:
        cur += 1.7
    ih = text_h(BIOMECH_INTRO, INNER_W, size=T_BODY, bulleted=False)
    if draw:
        text(slide, COL_L + PAD, cur + 0.1, INNER_W, ih + 0.3, BIOMECH_INTRO,
             size=T_BODY, color=INK, line=0.95)
    cur += 0.1 + ih + 0.4

    p2h = img_h("fig_p2_yaw.png", INNER_W)
    if draw:
        image(slide, "fig_p2_yaw.png", COL_L + PAD, cur, INNER_W, p2h + 0.01)
    cur += p2h + 0.15

    ch = text_h(YAW_CAPTION, INNER_W, size=T_BODY, bulleted=False)
    if draw:
        text(slide, COL_L + PAD, cur, INNER_W, ch + 0.2, YAW_CAPTION, size=T_BODY,
             color=INK, line=0.95)
    cur += ch + 0.6

    if draw:
        cur = sub(slide, COL_L + PAD, cur, INNER_W, "Estimator")
    else:
        cur += 1.7
    if draw:
        cw = (INNER_W - 1.2) / 3
        chip(slide, COL_L + PAD, cur + 0.1, cw, CHIP_H,
             [[("θ", {}), ("robot", {"sub": True}), ("   (left hip)", {})]])
        chip(slide, COL_L + PAD + cw + 0.6, cur + 0.1, cw, CHIP_H,
             [[("θ", {}), ("robot", {"sub": True}), ("   (right hip)", {})]])
        chip(slide, COL_L + PAD + 2 * (cw + 0.6), cur + 0.1, cw, CHIP_H,
             [[("θ̃", {}), ("yaw", {"sub": True}),
               ("   trunk-IMU yaw\n(0.1 Hz high-pass)", {})]])
    cur += 0.1 + CHIP_H + 0.3

    calib_h = text_h(CALIB_NOTE, INNER_W, size=T_BODY, bulleted=False)
    if draw:
        text(slide, COL_L + PAD, cur, INNER_W, calib_h + 0.2, CALIB_NOTE,
             size=T_BODY, color=INK, line=0.95)
    cur += calib_h + 0.6

    arch_w = INNER_W * 0.66
    arch_h = img_h("fig_p_architecture.png", arch_w)
    if draw:
        # box width = full INNER_W (wider than the image's own arch_w) so the
        # image, which stays exactly arch_w wide, is horizontally centred in
        # the column — matching the centred lines below it.
        image(slide, "fig_p_architecture.png", COL_L + PAD, cur, INNER_W, arch_h + 0.01)
    cur += arch_h + 0.35

    param_h = text_h(PARAM_LINE, INNER_W, size=T_BODY, bulleted=False)
    if draw:
        text(slide, COL_L + PAD, cur, INNER_W, param_h + 0.2, PARAM_LINE, size=T_BODY,
             bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    cur += param_h + 0.25

    close_h = text_h(METHOD_CLOSE, INNER_W, size=T_BODY, bulleted=False)
    if draw:
        text(slide, COL_L + PAD, cur, INNER_W, close_h + 0.2, METHOD_CLOSE,
             size=T_BODY, color=INK, align=PP_ALIGN.CENTER)
    cur += close_h
    return cur


SEC_METHOD = SEC_HEAD_H + (layout_methodology(0, draw=False)) + BOTTOM_PAD

# ---------------------------------------------------------- Experiment layout
EXP_ITEMS = [
    "16 healthy adults · 11.7 h of robot–mocap data at 100 Hz",
    "0.75 / 1.00 / 1.25 m/s and stop-and-go; assist lv0 / lv4 / lv7",
    "Reference: 12-camera Vicon, three thigh markers, T-pose zeroed",
    "Leave-one-subject-out over participants",
    "Motion capture trains and scores the estimate; never used at "
    "inference.",
]
EXP_IMG_W = 20.5
EXP_SIZE = T_BODY
EXP_TEXT_X = COL_R + TEXT_X_OFFSET
EXP_TEXT_W = ROW_TEXT_W
EXP_ROW_H = max(img_h("fig_p4_setup.png", EXP_IMG_W),
                text_h(EXP_ITEMS, EXP_TEXT_W, size=EXP_SIZE, space=8))
SEC_EXPERIMENT = SEC_HEAD_H + EXP_ROW_H + BOTTOM_PAD

# -------------------------------------------------------------- Results layout
RESULT_BLOCKS = [
    ("Accuracy on unseen participants", "fig_r1_indomain.png", [
        "16-fold leave-one-subject-out: −3.15° vs the raw robot angle "
        "(95% CI −3.79 to −2.47°, p < 0.001), lower in all 16 held-out "
        "participants.",
        "−0.28° vs the no-yaw network (p = 0.039); −0.70° vs linear "
        "Ridge (p = 0.002).",
    ]),
    ("Conditions the model never saw", "fig_r2_ood.png", [
        "Each target condition (assistance, faster walking, stop-and-go) was "
        "excluded from training, validation and model selection.",
        "Overall MAE fell 35–53% relative to the raw robot angle across all "
        "three shifts.",
    ]),
    ("Where the raw angle fails most", "fig_r3_assist.png", [
        "High-flexion MAE: 14.23° → 4.21° (−70.4%) under assistance.",
        "Raw error grows with assistance level (7.67° → 8.94°); Proposed "
        "stays flat (3.74° → 4.09°).",
    ]),
    ("What this makes possible", "fig_r4_application.png", [
        "Raw overstates the lv0→lv7 ROM change by 172%, Proposed only 50%; "
        "samples above 5° error fall from 58.4% to 30.1%.",
        "2.642 ms mean latency and 8,658 parameters on a Jetson Orin NX may "
        "enable estimator-in-the-loop use; closed-loop control was not "
        "verified.",
    ]),
]
block_h = [1.7 + RES_GAP_TITLE_IMG + chart_h(f) + RES_GAP_IMG_TXT + text_h(b, INNER_W)
           for _, f, b in RESULT_BLOCKS]
SEC_RESULTS = (SEC_HEAD_H + sum(block_h) + RES_GAP_BLOCK * (len(RESULT_BLOCKS) - 1)
               + BOTTOM_PAD)

print(f"left column total  = {SEC_BACKGROUND + GAP + SEC_ROBOT + GAP + SEC_METHOD:.2f} cm")
print(f"right column total = {SEC_EXPERIMENT + GAP + SEC_RESULTS:.2f} cm")
print(f"budget (BODY_Y..FOOT_Y) = {FOOT_Y - BODY_Y:.2f} cm")

# ============================================================= LEFT COLUMN ====
y = BODY_Y

iy = section(slide, COL_L, y, COL_W, "Background")
image(slide, "fig_p1_concept.png", COL_L + PAD, iy, IMG_SLOT_W, BG_ROW_H)
bullets(slide, BG_TEXT_X, iy, BG_TEXT_W, BG_ITEMS, size=BG_SIZE, space=11,
        h=BG_ROW_H, anchor=MSO_ANCHOR.MIDDLE)
y += SEC_BACKGROUND + GAP

iy = section(slide, COL_L, y, COL_W, "Robot System")
image(slide, "fig_p3_sensors.png", COL_L + PAD, iy, IMG_SLOT_W, RS_ROW_H)
bullets(slide, RS_TEXT_X, iy, RS_TEXT_W, RS_ITEMS, size=RS_SIZE, space=9,
        h=RS_ROW_H, anchor=MSO_ANCHOR.MIDDLE)
y += SEC_ROBOT + GAP

iy = section(slide, COL_L, y, COL_W, "Methodology")
layout_methodology(iy, draw=True)

# ============================================================ RIGHT COLUMN ====
y = BODY_Y

iy = section(slide, COL_R, y, COL_W, "Experiment")
image(slide, "fig_p4_setup.png", COL_R + PAD, iy, IMG_SLOT_W, EXP_ROW_H)
bullets(slide, EXP_TEXT_X, iy, EXP_TEXT_W, EXP_ITEMS, size=EXP_SIZE, space=8,
        h=EXP_ROW_H, anchor=MSO_ANCHOR.MIDDLE)
y += SEC_EXPERIMENT + GAP

iy = section(slide, COL_R, y, COL_W, "Results")
cur = iy
for (title, fname, items), bh in zip(RESULT_BLOCKS, block_h):
    cur = sub(slide, COL_R + PAD, cur, INNER_W, "•  " + title)
    cur += RES_GAP_TITLE_IMG
    ih = chart_h(fname)
    image(slide, fname, COL_R + PAD, cur, INNER_W, ih + 0.01)
    cur += ih + RES_GAP_IMG_TXT
    bullets(slide, COL_R + PAD, cur, INNER_W, items)
    cur += text_h(items, INNER_W) + RES_GAP_BLOCK
y += SEC_RESULTS

# ================================================================== footer ====
rect(slide, MARGIN, FOOT_Y, CONTENT_W, 0.05, fill=CARD_EDGE)
text(slide, MARGIN, FOOT_Y + 0.35, CONTENT_W, 1.6,
     "Evaluated in 16 healthy adult male participants, one hip exoskeleton and one "
     "low-magnitude assistance profile.   ·   KAIST IRB 2025-295   ·   "
     "Supported by NRF Korea (MSIT) No. 2022R1A3B1077880",
     size=19, color=INK, align=PP_ALIGN.CENTER, line=0.95)

prs.save(OUT)
print(f"saved {OUT}  ({PAGE_W} x {PAGE_H} cm, {len(slide.shapes)} objects)")
